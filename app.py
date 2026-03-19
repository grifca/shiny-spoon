"""
Campaign Performance Reporting App
"""

import streamlit as st
import pandas as pd
import numpy as np
from io import BytesIO
import os
import sys
import json

# ── Persistent config ──────────────────────────────────────────────────────────
_CONFIG_FILE = os.path.join(os.path.dirname(__file__), ".app_config.json")

def _load_config():
    if os.path.exists(_CONFIG_FILE):
        try:
            with open(_CONFIG_FILE) as f:
                return json.load(f)
        except Exception:
            pass
    return {}

def _save_config(data: dict):
    try:
        existing = _load_config()
        existing.update(data)
        with open(_CONFIG_FILE, "w") as f:
            json.dump(existing, f)
    except Exception:
        pass

# Seed session state from persisted config on first load
if "config_loaded" not in st.session_state:
    _cfg = _load_config()
    if _cfg.get("anthropic_api_key"):
        st.session_state["anthropic_api_key"] = _cfg["anthropic_api_key"]
    st.session_state["config_loaded"] = True

# ── Page config ───────────────────────────────────────────────────────────────
st.set_page_config(
    page_title="Automated Performance",
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="expanded",
)

# ── Styling ───────────────────────────────────────────────────────────────────
st.markdown("""
<style>
  @import url('https://fonts.googleapis.com/css2?family=Nunito+Sans:wght@300;400;600;700;800&display=swap');

  html, body, [class*="css"] {
    font-family: 'Nunito Sans', 'Avenir', 'Helvetica Neue', sans-serif;
  }

  /* Hide Streamlit branding */
  #MainMenu {visibility: hidden;}
  footer {visibility: hidden;}
  header {visibility: hidden;}

  /* App background */
  .stApp { background-color: #F7F8FA; }

  /* Sidebar */
  [data-testid="stSidebar"] {
    background: #1A2B4A;
  }
  [data-testid="stSidebar"] * {
    color: #FFFFFF !important;
  }
  [data-testid="stSidebar"] label,
  [data-testid="stSidebar"] .stSelectbox label,
  [data-testid="stSidebar"] .stMultiSelect label,
  [data-testid="stSidebar"] .stFileUploader label,
  [data-testid="stSidebar"] .stTextInput label,
  [data-testid="stSidebar"] .stTextArea label,
  [data-testid="stSidebar"] .stRadio label,
  [data-testid="stSidebar"] .stDateInput label,
  [data-testid="stSidebar"] p {
    color: #FFFFFF !important;
    font-weight: 600;
  }

  /* Main area — all form labels must be dark and readable */
  .main label,
  .main .stSelectbox label,
  .main .stMultiSelect label,
  .main .stTextInput label,
  .main .stTextArea label,
  .main .stDateInput label,
  .main .stRadio label,
  .main .stCheckbox label,
  .main .stNumberInput label,
  .main [data-testid="stWidgetLabel"],
  .main [data-testid="stWidgetLabel"] p,
  .main [data-testid="stWidgetLabel"] span,
  div[data-testid="stMainBlockContainer"] [data-testid="stWidgetLabel"],
  div[data-testid="stMainBlockContainer"] [data-testid="stWidgetLabel"] p,
  div[data-testid="stMainBlockContainer"] label {
    color: #1A2B4A !important;
    font-weight: 700 !important;
    font-size: 0.8rem !important;
  }

  /* Selectbox/input text values in main area */
  .main .stSelectbox > div > div,
  .main .stMultiSelect > div > div {
    color: #1A2B4A !important;
  }

  /* Header bar */
  .app-header {
    background: #1A2B4A;
    padding: 1.5rem 2rem;
    margin: -1rem -1rem 2rem -1rem;
    border-bottom: 4px solid #DF552C;
  }
  .app-header h1 {
    color: #FFFFFF;
    font-size: 1.6rem;
    font-weight: 800;
    margin: 0;
    letter-spacing: -0.02em;
  }
  .app-header p {
    color: #8C9DB5;
    font-size: 0.85rem;
    margin: 0.25rem 0 0 0;
  }

  /* KPI cards */
  .kpi-card {
    background: #FFFFFF;
    border-radius: 8px;
    padding: 1.25rem 1.5rem;
    box-shadow: 0 1px 4px rgba(0,0,0,0.08);
    border-top: 3px solid #1A2B4A;
  }
  .kpi-card.accent { border-top-color: #DF552C; }
  .kpi-value {
    font-size: 2rem;
    font-weight: 800;
    color: #1A2B4A;
    line-height: 1;
  }
  .kpi-value.magenta { color: #DF552C; }
  .kpi-label {
    font-size: 0.75rem;
    font-weight: 600;
    color: #8C8C8C;
    text-transform: uppercase;
    letter-spacing: 0.05em;
    margin-top: 0.4rem;
  }
  .kpi-delta {
    font-size: 0.8rem;
    color: #22C55E;
    font-weight: 600;
    margin-top: 0.2rem;
  }
  .kpi-delta.down { color: #EF4444; }

  /* Section headers */
  .section-header {
    font-size: 0.7rem;
    font-weight: 700;
    color: #8C8C8C;
    text-transform: uppercase;
    letter-spacing: 0.1em;
    padding-bottom: 0.5rem;
    border-bottom: 2px solid #DF552C;
    margin-bottom: 1rem;
  }

  /* Insight cards */
  .insight-card {
    background: #1A2B4A;
    border-radius: 8px;
    padding: 1rem 1.25rem;
    color: white;
    margin-bottom: 0.75rem;
  }
  .insight-card .insight-title {
    font-size: 0.7rem;
    font-weight: 700;
    color: #DF552C;
    text-transform: uppercase;
    letter-spacing: 0.08em;
  }
  .insight-card .insight-text {
    font-size: 0.9rem;
    font-weight: 600;
    color: #FFFFFF;
    margin-top: 0.25rem;
    line-height: 1.4;
  }

  /* Table styling */
  .styled-table {
    width: 100%;
    border-collapse: collapse;
    font-size: 0.85rem;
  }
  .styled-table th {
    background: #1A2B4A;
    color: white;
    padding: 0.6rem 0.75rem;
    text-align: right;
    font-size: 0.7rem;
    font-weight: 700;
    text-transform: uppercase;
    letter-spacing: 0.05em;
  }
  .styled-table th:first-child { text-align: left; }
  .styled-table td {
    padding: 0.55rem 0.75rem;
    border-bottom: 1px solid #F0F0F0;
    color: #333333;
    text-align: right;
  }
  .styled-table td:first-child { text-align: left; }
  .styled-table tr:nth-child(even) td { background: #F9FAFB; }
  .styled-table tr:hover td { background: #EEF2FF; }
  .highlight-row td { background: #FFF0F6 !important; font-weight: 700; }

  /* Buttons */
  .stButton > button {
    background: #DF552C !important;
    color: white !important;
    border: none !important;
    border-radius: 6px !important;
    font-weight: 700 !important;
    font-family: 'Nunito Sans', 'Avenir', sans-serif !important;
    padding: 0.5rem 1.5rem !important;
    letter-spacing: 0.03em;
  }
  .stButton > button:hover {
    background: #C4105E !important;
    transform: translateY(-1px);
    box-shadow: 0 4px 12px rgba(232, 31, 118, 0.3) !important;
  }

  /* Download button */
  .stDownloadButton > button {
    background: #1A2B4A !important;
    color: white !important;
    border: none !important;
    border-radius: 6px !important;
    font-weight: 700 !important;
    width: 100%;
  }

  /* Tabs */
  .stTabs [data-baseweb="tab-list"] {
    gap: 0;
    border-bottom: 2px solid #E0E0E0;
  }
  .stTabs [data-baseweb="tab"] {
    font-weight: 600;
    font-size: 0.85rem;
    color: #8C8C8C;
    padding: 0.6rem 1.2rem;
    border-bottom: 2px solid transparent;
  }
  .stTabs [aria-selected="true"] {
    color: #1A2B4A !important;
    border-bottom: 2px solid #DF552C !important;
    background: transparent !important;
  }

  /* Status indicators */
  .status-good { color: #22C55E; font-weight: 700; }
  .status-warn { color: #F59E0B; font-weight: 700; }
  .status-bad  { color: #EF4444; font-weight: 700; }

  /* Connection status */
  .conn-badge {
    display: inline-block;
    padding: 0.2rem 0.6rem;
    border-radius: 20px;
    font-size: 0.7rem;
    font-weight: 700;
    text-transform: uppercase;
    letter-spacing: 0.05em;
  }
  .conn-badge.connected { background: #DCFCE7; color: #166534; }
  .conn-badge.disconnected { background: #FEE2E2; color: #991B1B; }

  /* Progress bar override */
  .stProgress > div > div { background-color: #DF552C !important; }

</style>
""", unsafe_allow_html=True)


# ── Helpers ───────────────────────────────────────────────────────────────────

def fmt_currency(v, decimals=0):
    if v is None or (isinstance(v, float) and np.isnan(v)):
        return "—"
    if v >= 1_000_000:
        return f"${v/1_000_000:.2f}MM"
    if v >= 1_000:
        return f"${v/1_000:.0f}K"
    return f"${v:,.{decimals}f}"

def fmt_pct(v, decimals=2):
    if v is None or (isinstance(v, float) and np.isnan(v)):
        return "—"
    return f"{v*100:.{decimals}f}%"

def fmt_int(v):
    if v is None or (isinstance(v, float) and np.isnan(v)):
        return "—"
    v = int(v)
    if v >= 1_000_000:
        return f"{v/1_000_000:.1f}MM"
    if v >= 1_000:
        return f"{v/1_000:.0f}K"
    return f"{v:,}"

def fmt_table_int(v):
    """Comma-separated integer for data tables. No decimals."""
    if v is None or (isinstance(v, float) and np.isnan(v)):
        return "—"
    return f"{int(round(v)):,}"

def fmt_table_currency(v):
    """Comma-separated dollar amount for data tables. Decimals only if < 10."""
    if v is None or (isinstance(v, float) and np.isnan(v)):
        return "—"
    if abs(v) < 10:
        return f"${v:,.2f}"
    return f"${int(round(v)):,}"

def safe_div(n, d, scale=1.0):
    if d and d != 0:
        return (n or 0) / d * scale
    return None

def compute_kpis(df):
    spend       = df["Spend"].sum()
    impressions = df["Impressions"].sum()
    clicks      = df["Clicks"].sum()
    visits      = df["Site_Visits"].sum()
    conversions = df["Conversion"].sum()
    return {
        "spend": spend,
        "impressions": impressions,
        "clicks": clicks,
        "visits": visits,
        "conversions": conversions,
        "ctr": safe_div(clicks, impressions),
        "cpc": safe_div(spend, clicks),
        "cpm": safe_div(spend, impressions, 1000),
        "cvr": safe_div(conversions, clicks),
        "cpa": safe_div(spend, conversions),
        "site_visit_rate": safe_div(visits, clicks),
    }

def generate_insights(df, kpis):
    insights = []
    # Channel analysis
    if "Channel" in df.columns:
        ch = df.groupby("Channel").agg(
            Spend=("Spend","sum"), Conversions=("Conversion","sum"), Clicks=("Clicks","sum")
        ).reset_index()
        ch = ch[ch["Spend"] > 0].copy()
        if not ch.empty:
            top_conv = ch.sort_values("Conversions", ascending=False).iloc[0]
            top_eff = ch[ch["Conversions"] > 0].copy()
            if not top_eff.empty:
                top_eff["CPA"] = top_eff["Spend"] / top_eff["Conversions"]
                best_cpa = top_eff.sort_values("CPA").iloc[0]
                insights.append({
                    "title": "Top Converting Channel",
                    "text": f"{top_conv['Channel']} drives the most conversions ({fmt_int(top_conv['Conversions'])}). "
                            f"Best CPA: {best_cpa['Channel']} at {fmt_currency(best_cpa['CPA'], 2)}."
                })
            zero_conv = ch[(ch["Spend"]/ch["Spend"].sum() >= 0.05) & (ch["Conversions"] == 0)]
            if not zero_conv.empty:
                names = ", ".join(zero_conv["Channel"].tolist())
                insights.append({
                    "title": "⚠ Zero Conversion Channels",
                    "text": f"{names} account for ≥5% of spend but show zero conversions. Review attribution or reallocate budget."
                })

    # Platform analysis
    if "Platform" in df.columns:
        pl = df.groupby("Platform").agg(
            Spend=("Spend","sum"), Conversions=("Conversion","sum")
        ).reset_index()
        top_platform = pl.sort_values("Spend", ascending=False).iloc[0]
        insights.append({
            "title": "Dominant Platform",
            "text": f"{top_platform['Platform']} leads with {fmt_currency(top_platform['Spend'])} spend "
                    f"({top_platform['Spend']/pl['Spend'].sum()*100:.0f}% of total)."
        })

    # CPA insight
    if kpis["cpa"]:
        insights.append({
            "title": "Overall Efficiency",
            "text": f"Average CPA: {fmt_currency(kpis['cpa'], 2)} | CTR: {fmt_pct(kpis['ctr'])} | CVR: {fmt_pct(kpis['cvr'])}"
        })

    return insights[:4]


# ── Data loading ──────────────────────────────────────────────────────────────

@st.cache_data(show_spinner=False)
def load_excel(file_bytes, filename):
    xl = pd.read_excel(BytesIO(file_bytes), sheet_name=None)
    return xl


def try_redshift(host, port, db, user, password, query):
    """Connect to Redshift and execute a query. Returns (df, error)."""
    try:
        import psycopg2
        conn = psycopg2.connect(
            host=host, port=int(port), dbname=db, user=user, password=password,
            connect_timeout=10
        )
        df = pd.read_sql(query, conn)
        conn.close()
        return df, None
    except ImportError:
        return None, "psycopg2 not installed. Run: pip install psycopg2-binary"
    except Exception as e:
        return None, str(e)


# ── Deck generation ───────────────────────────────────────────────────────────

def _fig_to_bytes(fig, width=1200, height=380):
    """Export a plotly figure to PNG bytes."""
    try:
        import plotly.io as pio
        return pio.to_image(fig, format="png", width=width, height=height, scale=2)
    except Exception:
        return None

def _make_chart_common():
    return dict(
        plot_bgcolor="white", paper_bgcolor="white",
        font=dict(family="Avenir, Calibri", color="#333333"),
        title_font_color="#1A2B4A",
        margin=dict(l=80, r=100, t=60, b=80),
        uniformtext=dict(mode="hide", minsize=9),
    )

def generate_deck(df, config):
    """Generate PPTX deck from dataframe + config. Returns BytesIO."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Pt
        import math
        import plotly.express as px
        import plotly.graph_objects as go
        import io
    except ImportError:
        return None, "python-pptx not installed. Run: pip install python-pptx"

    NAVY    = RGBColor(0x1A, 0x2B, 0x4A)
    MAGENTA = RGBColor(0xDF, 0x55, 0x2C)
    WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
    LGRAY   = RGBColor(0xF4, 0xF4, 0xF4)
    CHARCOAL= RGBColor(0x33, 0x33, 0x33)
    MGRAY   = RGBColor(0x8C, 0x8C, 0x8C)

    W = Inches(13.33)
    H = Inches(7.5)

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    blank_layout = prs.slide_layouts[6]

    def add_rect(slide, x, y, w, h, fill=None, line=None, line_width=Pt(0)):
        from pptx.util import Pt
        shape = slide.shapes.add_shape(1, x, y, w, h)
        shape.line.fill.background()
        if fill:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill
        else:
            shape.fill.background()
        if line:
            shape.line.color.rgb = line
            shape.line.width = line_width
        else:
            shape.line.fill.background()
        return shape

    def add_text(slide, text, x, y, w, h, font_size=Pt(14), bold=False,
                 color=None, align=PP_ALIGN.LEFT, italic=False, wrap=True):
        txBox = slide.shapes.add_textbox(x, y, w, h)
        tf = txBox.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = font_size
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color or CHARCOAL
        try:
            run.font.name = "Avenir"
        except Exception:
            run.font.name = "Calibri"
        return txBox

    def header_bar(slide, title):
        bar = add_rect(slide, 0, 0, W, Inches(0.45), fill=NAVY)
        add_text(slide, title.upper(),
                 Inches(0.4), Inches(0.08), Inches(12), Inches(0.3),
                 font_size=Pt(11), bold=True, color=WHITE)

    kpis = compute_kpis(df)
    lob  = config.get("lob", "All")
    client = config.get("client", "Client")
    date_range = config.get("date_range", "")

    # ── SLIDE 1: Cover ────────────────────────────────────────────────────────
    s1 = prs.slides.add_slide(blank_layout)
    add_rect(s1, 0, 0, W, H, fill=NAVY)
    add_rect(s1, 0, H - Inches(0.08), W, Inches(0.08), fill=MAGENTA)
    add_text(s1, client.upper(), Inches(0.5), Inches(0.4), Inches(5), Inches(0.35),
             font_size=Pt(12), bold=True, color=MAGENTA)
    add_text(s1, lob, Inches(0.5), Inches(2.8), Inches(12), Inches(1.2),
             font_size=Pt(44), bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_text(s1, "Campaign Performance Summary",
             Inches(0.5), Inches(4.0), Inches(10), Inches(0.6),
             font_size=Pt(22), bold=False, color=WHITE)
    add_text(s1, date_range or "Full Campaign Flight",
             Inches(0.5), Inches(4.65), Inches(8), Inches(0.4),
             font_size=Pt(16), bold=False, color=MGRAY)
    add_text(s1, date_range or "",
             W - Inches(3.2), H - Inches(0.4), Inches(2.8), Inches(0.3),
             font_size=Pt(9), bold=False, color=MGRAY, align=PP_ALIGN.RIGHT)

    # ── SLIDE 2: KPI Summary + Monthly Trend (combined) ───────────────────────
    s2 = prs.slides.add_slide(blank_layout)
    add_rect(s2, 0, 0, W, H, fill=WHITE)
    header_bar(s2, "Monthly Performance Trend")

    # Layout constants
    card_gap  = Inches(0.18)
    card_w    = (W - Inches(0.4) * 2 - card_gap * 3) / 4
    row1_h    = Inches(1.35)
    row2_h    = Inches(0.82)
    row1_y    = Inches(0.62)
    row2_y    = row1_y + row1_h + Inches(0.14)
    chart_y   = row2_y + row2_h + Inches(0.16)
    start_x   = Inches(0.4)

    # Row 1: primary KPI cards
    kpi_items = [
        (fmt_currency(kpis["spend"]),  "Total Spend",      NAVY),
        (fmt_int(kpis["impressions"]), "Impressions",      NAVY),
        (fmt_int(kpis["conversions"]), "Conversions",      MAGENTA),
        (fmt_currency(kpis["cpa"], 2) if kpis["cpa"] else "—", "Avg CPA", MAGENTA),
    ]
    for i, (val, label, color) in enumerate(kpi_items):
        cx = start_x + i * (card_w + card_gap)
        add_rect(s2, cx, row1_y, card_w, row1_h, fill=LGRAY)
        add_rect(s2, cx, row1_y, Inches(0.055), row1_h, fill=color)
        add_text(s2, val,   cx + Inches(0.16), row1_y + Inches(0.18),
                 card_w - Inches(0.2), Inches(0.72),
                 font_size=Pt(30), bold=True, color=color)
        add_text(s2, label, cx + Inches(0.16), row1_y + Inches(0.95),
                 card_w - Inches(0.2), Inches(0.3),
                 font_size=Pt(10), bold=False, color=MGRAY)

    # Row 2: efficiency KPI cards
    eff_items = [
        (fmt_pct(kpis["ctr"]),  "CTR"),
        (fmt_pct(kpis["conv_rate"]) if kpis.get("conv_rate") else "—", "CVR"),
        (fmt_currency(kpis["cpc"], 2) if kpis["cpc"] else "—", "CPC"),
        (fmt_currency(kpis["cpm"], 2) if kpis["cpm"] else "—", "CPM"),
    ]
    for i, (val, label) in enumerate(eff_items):
        ex = start_x + i * (card_w + card_gap)
        add_rect(s2, ex, row2_y, card_w, row2_h, fill=LGRAY)
        add_text(s2, val,   ex + Inches(0.14), row2_y + Inches(0.1),
                 card_w - Inches(0.18), Inches(0.45),
                 font_size=Pt(20), bold=True, color=NAVY)
        add_text(s2, label, ex + Inches(0.14), row2_y + Inches(0.55),
                 card_w - Inches(0.18), Inches(0.22),
                 font_size=Pt(9), bold=False, color=MGRAY)

    # Monthly chart embedded below cards
    if "Date" in df.columns:
        df_t2 = df.copy()
        df_t2["Date"] = pd.to_datetime(df_t2["Date"], errors="coerce")
        monthly2 = df_t2.groupby(df_t2["Date"].dt.to_period("M")).agg(
            Spend=("Spend", "sum"), Conversions=("Conversion", "sum")
        ).reset_index()
        monthly2["Month"] = monthly2["Date"].astype(str)
        fig_kpi_trend = go.Figure()
        fig_kpi_trend.add_trace(go.Bar(
            x=monthly2["Month"], y=monthly2["Spend"],
            name="Spend", marker_color="#1A2B4A", yaxis="y1",
        ))
        fig_kpi_trend.add_trace(go.Scatter(
            x=monthly2["Month"], y=monthly2["Conversions"],
            name="Conversions", line=dict(color="#DF552C", width=3),
            mode="lines+markers", marker=dict(size=7), yaxis="y2",
        ))
        fig_kpi_trend.update_layout(
            **_make_chart_common(),
            title="Monthly Spend vs. Conversions",
            xaxis=dict(tickfont=dict(color="#333333"), tickangle=-30, showgrid=False),
            yaxis=dict(
                title=dict(text="Spend ($)", font=dict(color="#333333")),
                tickfont=dict(color="#333333"), showgrid=True, gridcolor="#EEEEEE",
                tickprefix="$", tickformat=",.0f",
            ),
            yaxis2=dict(
                title=dict(text="Conversions", font=dict(color="#DF552C")),
                tickfont=dict(color="#DF552C"), overlaying="y", side="right",
            ),
            legend=dict(orientation="h", yanchor="bottom", y=1.02,
                        font=dict(color="#333333")),
            height=420,
        )
        CHART_W_PX, CHART_H_PX = 1400, 420
        chart_bytes = _fig_to_bytes(fig_kpi_trend, width=CHART_W_PX, height=CHART_H_PX)
        if chart_bytes:
            avail_w = W - Inches(0.5)
            avail_h = H - chart_y - Inches(0.1)
            aspect   = CHART_W_PX / CHART_H_PX
            fit_w    = avail_w
            fit_h    = fit_w / aspect
            if fit_h > avail_h:
                fit_h = avail_h
                fit_w = fit_h * aspect
            left = (W - fit_w) / 2
            s2.shapes.add_picture(io.BytesIO(chart_bytes), left, chart_y, fit_w, fit_h)

    def add_chart_slide(title, fig, notes=None):
        s = prs.slides.add_slide(blank_layout)
        add_rect(s, 0, 0, W, H, fill=WHITE)
        header_bar(s, title)
        IMG_W, IMG_H = 1600, 480
        img_bytes = _fig_to_bytes(fig, width=IMG_W, height=IMG_H)
        if img_bytes:
            # Max available area (leave margin + room for notes)
            max_w = W - Inches(0.6)
            max_h = H - Inches(0.55) - Inches(1.0)  # header + notes space
            # Scale to fit while preserving aspect ratio
            aspect = IMG_W / IMG_H
            fit_w = max_w
            fit_h = fit_w / aspect
            if fit_h > max_h:
                fit_h = max_h
                fit_w = fit_h * aspect
            # Center horizontally
            left = (W - fit_w) / 2
            top = Inches(0.6)
            s.shapes.add_picture(io.BytesIO(img_bytes), left, top, fit_w, fit_h)
        if notes:
            add_text(s, notes, Inches(0.4), H - Inches(0.75), W - Inches(0.8), Inches(0.6),
                     font_size=Pt(10), bold=False, color=MGRAY, italic=True)
        return s

    # ── SLIDE 3: Channel breakdown ────────────────────────────────────────────
    if "Channel" in df.columns:
        ch_df = df.groupby("Channel").agg(
            Spend=("Spend","sum"), Conversions=("Conversion","sum"), Clicks=("Clicks","sum")
        ).reset_index().sort_values("Spend", ascending=False).head(8)
        fig_ch = px.bar(ch_df, x="Spend", y="Channel", orientation="h",
                        color_discrete_sequence=["#1A2B4A"], title="Spend by Channel",
                        text=ch_df["Spend"].apply(fmt_currency))
        fig_ch.update_layout(**_make_chart_common(), height=400,
                             xaxis=dict(tickfont=dict(color="#333333"), showgrid=True, gridcolor="#EEEEEE"),
                             yaxis=dict(tickfont=dict(color="#333333"), automargin=True))
        fig_ch.update_traces(marker_line_width=0, textposition="auto",
                             textfont=dict(color="white", size=11), cliponaxis=False)
        add_chart_slide("Channel Performance", fig_ch,
                        notes="Sorted by spend. Conversion counts labelled on bars.")

    # ── SLIDE 4: Platform breakdown ───────────────────────────────────────────
    if "Platform" in df.columns:
        pl_df = df.groupby("Platform").agg(
            Spend=("Spend","sum"), Conversions=("Conversion","sum"),
            Clicks=("Clicks","sum"), Impressions=("Impressions","sum")
        ).reset_index().sort_values("Spend", ascending=False).head(6)
        pl_df["CPA"] = pl_df.apply(
            lambda r: r["Spend"]/r["Conversions"] if r["Conversions"] > 0 else None, axis=1)
        fig_pl = px.bar(pl_df, x="Platform", y="Spend",
                        color_discrete_sequence=["#1A2B4A"], title="Spend by Platform",
                        text=pl_df["Spend"].apply(fmt_currency))
        fig_pl.update_layout(**_make_chart_common(), height=400,
                             xaxis=dict(tickfont=dict(color="#333333"), automargin=True),
                             yaxis=dict(title=dict(text="Spend ($)", font=dict(color="#333333")),
                                        tickfont=dict(color="#333333"), showgrid=True, gridcolor="#EEEEEE"))
        fig_pl.update_traces(marker_line_width=0, textposition="auto",
                             textfont=dict(color="white", size=11), cliponaxis=False)
        add_chart_slide("Platform Performance", fig_pl)

    # ── SLIDE 5: Product CPA chart ────────────────────────────────────────────
    prod_col_deck = next((c for c in ["Product","LOB","Campaign_Initiative"] if c in df.columns), None)
    if prod_col_deck:
        pr_d = df.groupby(prod_col_deck).agg(
            Spend=("Spend","sum"), Conversions=("Conversion","sum")
        ).reset_index()
        pr_d["CPA"] = pr_d.apply(lambda r: r["Spend"]/r["Conversions"] if r["Conversions"]>0 else None, axis=1)
        pr_d = pr_d.dropna(subset=["CPA"]).sort_values("CPA").head(10)
        fig_pr = px.bar(pr_d, x=prod_col_deck, y="CPA",
                        color_discrete_sequence=["#DF552C"], title="CPA by Product / Segment",
                        text=pr_d["CPA"].apply(lambda v: fmt_table_currency(v)))
        fig_pr.update_layout(**_make_chart_common(), height=400,
                             xaxis=dict(tickfont=dict(color="#333333"), automargin=True,
                                        tickangle=-30),
                             yaxis=dict(title=dict(text="CPA ($)", font=dict(color="#333333")),
                                        tickfont=dict(color="#333333"), showgrid=True, gridcolor="#EEEEEE"))
        fig_pr.update_traces(marker_line_width=0, textposition="auto",
                             textfont=dict(color="white", size=11), cliponaxis=False)
        add_chart_slide("Product Performance", fig_pr, notes="Sorted by CPA ascending (best to worst).")

    # Monthly trend is now embedded in slide 2 (KPI + trend combined slide)

    # ── SLIDE 7: Insights ─────────────────────────────────────────────────────
    s5 = prs.slides.add_slide(blank_layout)
    add_rect(s5, 0, 0, W, H, fill=LGRAY)
    add_rect(s5, 0, 0, W, Inches(0.45), fill=MAGENTA)
    add_text(s5, "KEY INSIGHTS",
             Inches(0.4), Inches(0.08), Inches(12), Inches(0.3),
             font_size=Pt(11), bold=True, color=WHITE)

    insights = generate_insights(df, kpis)
    card_h_ins = Inches(1.4)
    card_gap = Inches(0.2)
    cols = 2
    ins_w = (W - Inches(1.0) - Inches(0.2)) / cols
    start_y_ins = Inches(0.65)

    for idx, ins in enumerate(insights):
        col = idx % cols
        row = idx // cols
        cx = Inches(0.5) + col * (ins_w + Inches(0.2))
        cy = start_y_ins + row * (card_h_ins + card_gap)
        add_rect(s5, cx, cy, ins_w, card_h_ins, fill=NAVY)
        add_text(s5, ins["title"].upper(),
                 cx + Inches(0.2), cy + Inches(0.15),
                 ins_w - Inches(0.4), Inches(0.3),
                 font_size=Pt(8), bold=True, color=MAGENTA)
        add_text(s5, ins["text"],
                 cx + Inches(0.2), cy + Inches(0.45),
                 ins_w - Inches(0.4), Inches(0.85),
                 font_size=Pt(11), bold=False, color=WHITE)

    # Closing slide
    s6 = prs.slides.add_slide(blank_layout)
    add_rect(s6, 0, 0, W, H, fill=NAVY)
    add_rect(s6, 0, H - Inches(0.08), W, Inches(0.08), fill=MAGENTA)
    add_text(s6, "Campaign Analytics",
             Inches(0.5), Inches(2.5), Inches(12), Inches(0.8),
             font_size=Pt(36), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s6, client.upper() if client else "",
             Inches(0.5), Inches(3.4), Inches(12), Inches(0.5),
             font_size=Pt(18), bold=False, color=MGRAY, align=PP_ALIGN.CENTER)
    ds = f"Data: {lob} | {date_range}" if date_range else f"Data: {lob}"
    add_text(s6, ds,
             Inches(0.5), H - Inches(0.6), Inches(12), Inches(0.3),
             font_size=Pt(9), bold=False, color=MGRAY, align=PP_ALIGN.CENTER)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf, None


# ── Sidebar ───────────────────────────────────────────────────────────────────

with st.sidebar:
    st.markdown("""
    <div style="padding:0.5rem 0 1.5rem 0; border-bottom:1px solid rgba(255,255,255,0.15); margin-bottom:1.5rem;">
        <div style="font-size:1.2rem;font-weight:800;color:#FFFFFF;letter-spacing:-0.02em;">
            Automated Performance
        </div>
        <div style="font-size:0.75rem;color:#8C9DB5;margin-top:0.25rem;">
            MKS Reporting Automation Demo
        </div>
    </div>
    """, unsafe_allow_html=True)

    data_source = st.radio(
        "DATA SOURCE",
        ["Upload Excel", "SharePoint Link", "Redshift"],
        index=0,
    )

    st.markdown("---")
    st.markdown("**REPORT SETTINGS**")

    client_name = st.text_input("Client Name", value="", placeholder="Client name", key="client")
    report_title = st.text_input("Report Title", value="Campaign Performance Summary")


    df_raw = None
    redshift_connected = False

    # ── Disk cache paths ───────────────────────────────────────────────────────
    _CACHE_DIR  = os.path.join(os.path.dirname(__file__), ".data_cache")
    _CACHE_FILE = os.path.join(_CACHE_DIR, "cached_data.pkl")
    _META_FILE  = os.path.join(_CACHE_DIR, "cache_meta.json")
    os.makedirs(_CACHE_DIR, exist_ok=True)

    def _save_cache(df, source_name):
        import pickle, datetime
        df.to_pickle(_CACHE_FILE)
        with open(_META_FILE, "w") as f:
            json.dump({
                "source": source_name,
                "refreshed": datetime.datetime.now().isoformat(),
                "rows": len(df),
            }, f)

    def _load_cache():
        import pickle
        if os.path.exists(_CACHE_FILE) and os.path.exists(_META_FILE):
            try:
                df = pd.read_pickle(_CACHE_FILE)
                with open(_META_FILE) as f:
                    meta = json.load(f)
                return df, meta
            except Exception:
                pass
        return None, None

    if data_source == "Upload Excel":
        uploaded = st.file_uploader(
            "Drop Excel file here",
            type=["xlsx", "xls"],
            help="Supports multi-sheet workbooks",
        )
        if uploaded:
            with st.spinner("Reading file..."):
                sheets = load_excel(uploaded.read(), uploaded.name)
            sheet_names = list(sheets.keys())
            selected_sheet = st.selectbox("Sheet", sheet_names)
            df_raw = sheets[selected_sheet]
            _save_cache(df_raw, uploaded.name)
            st.success(f"✓ {len(df_raw):,} rows loaded")
        else:
            # Auto-reload from disk cache
            cached_df, cached_meta = _load_cache()
            if cached_df is not None:
                df_raw = cached_df
                import datetime
                refreshed_dt = datetime.datetime.fromisoformat(cached_meta["refreshed"])
                st.markdown(
                    f"<p style='font-size:0.72rem;color:#8C9DB5;margin-top:0.3rem;'>"
                    f"↻ Last loaded: <b>{cached_meta['source']}</b><br>"
                    f"{refreshed_dt.strftime('%b %-d, %Y at %-I:%M %p')} · {cached_meta['rows']:,} rows"
                    f"</p>",
                    unsafe_allow_html=True,
                )
                if st.button("Clear cached data", key="clear_cache"):
                    import shutil
                    shutil.rmtree(_CACHE_DIR, ignore_errors=True)
                    st.rerun()

    elif data_source == "SharePoint Link":
        st.markdown("**SHAREPOINT LINK**")
        sp_url = st.text_input(
            "Paste sharing URL",
            placeholder="https://yourorg.sharepoint.com/:x:/s/...",
            help="Must be an 'Anyone with the link' public share. Copy from SharePoint → Share → Anyone with the link.",
        )
        sp_sheet = st.text_input("Sheet name (optional)", placeholder="Leave blank for first sheet")

        def _sp_to_download_url(url):
            """Convert a SharePoint sharing URL to a direct download URL."""
            import urllib.parse
            # Handle /:x:/ or /:f:/ style sharing links
            if "/:x:/" in url or "/:X:/" in url or "download=1" in url:
                # Already a download or embed link — add download param
                sep = "&" if "?" in url else "?"
                return url + sep + "download=1"
            # For sharing links like https://org.sharepoint.com/:x:/s/site/Eabc123
            # Convert to the download endpoint
            parsed = urllib.parse.urlparse(url)
            # Encode the full URL as base64 for the _layouts/15/download.aspx endpoint
            import base64
            encoded = base64.b64encode(url.encode()).decode()
            # Remove padding
            encoded = encoded.rstrip("=")
            download_url = f"{parsed.scheme}://{parsed.netloc}/_layouts/15/download.aspx?share={encoded}"
            return download_url

        def _fetch_sharepoint(url, sheet_name=None):
            import requests, io
            try:
                dl_url = _sp_to_download_url(url)
                headers = {"User-Agent": "Mozilla/5.0"}
                resp = requests.get(dl_url, headers=headers, timeout=30, allow_redirects=True)
                if resp.status_code != 200:
                    return None, f"HTTP {resp.status_code} — check that the link is set to 'Anyone with the link'."
                content_type = resp.headers.get("Content-Type", "")
                if "html" in content_type.lower():
                    return None, "Received an HTML page instead of a file — the link requires Microsoft login (not publicly shared)."
                file_bytes = io.BytesIO(resp.content)
                sheets = load_excel(file_bytes.read(), "sharepoint.xlsx")
                if sheet_name and sheet_name in sheets:
                    return sheets[sheet_name], None
                return list(sheets.values())[0], None
            except Exception as e:
                return None, str(e)

        if st.button("Load from SharePoint", key="sp_load"):
            if sp_url:
                with st.spinner("Fetching file from SharePoint..."):
                    sp_df, sp_err = _fetch_sharepoint(sp_url, sp_sheet or None)
                if sp_err:
                    st.error(f"Could not load file: {sp_err}")
                else:
                    _save_cache(sp_df, sp_url.split("/")[-1] or "sharepoint_file.xlsx")
                    st.session_state["sp_df"] = sp_df
                    st.success(f"✓ {len(sp_df):,} rows loaded from SharePoint")
            else:
                st.warning("Paste a SharePoint sharing URL above.")

        if "sp_df" in st.session_state:
            df_raw = st.session_state["sp_df"]
        else:
            cached_df, cached_meta = _load_cache()
            if cached_df is not None:
                df_raw = cached_df
                import datetime
                refreshed_dt = datetime.datetime.fromisoformat(cached_meta["refreshed"])
                st.markdown(
                    f"<p style='font-size:0.72rem;color:#8C9DB5;margin-top:0.3rem;'>"
                    f"↻ Last loaded: <b>{cached_meta['source']}</b><br>"
                    f"{refreshed_dt.strftime('%b %-d, %Y at %-I:%M %p')} · {cached_meta['rows']:,} rows"
                    f"</p>",
                    unsafe_allow_html=True,
                )

    else:
        st.markdown("**REDSHIFT CONNECTION**")
        rs_host = st.text_input("Host", placeholder="your-cluster.region.redshift.amazonaws.com")
        rs_port = st.text_input("Port", value="5439")
        rs_db   = st.text_input("Database", placeholder="analytics")
        rs_user = st.text_input("Username")
        rs_pass = st.text_input("Password", type="password")
        rs_query = st.text_area(
            "SQL Query",
            value="SELECT * FROM campaign_data WHERE lob = 'Personal Risk Services' LIMIT 10000",
            height=100,
        )
        if st.button("Connect & Run Query"):
            if all([rs_host, rs_db, rs_user, rs_pass]):
                with st.spinner("Connecting to Redshift..."):
                    df_raw, err = try_redshift(rs_host, rs_port, rs_db, rs_user, rs_pass, rs_query)
                if err:
                    st.error(f"Connection failed: {err}")
                else:
                    st.session_state["redshift_df"] = df_raw
                    st.success(f"✓ {len(df_raw):,} rows returned")
            else:
                st.warning("Fill in all connection fields.")

        if "redshift_df" in st.session_state:
            df_raw = st.session_state["redshift_df"]
            redshift_connected = True
            st.markdown('<span class="conn-badge connected">● CONNECTED</span>', unsafe_allow_html=True)


# ── Main ──────────────────────────────────────────────────────────────────────

st.markdown("""
<div class="app-header">
    <h1>Automated Performance</h1>
    <p>Upload an Excel workbook or connect to Redshift to generate a performance summary deck</p>
</div>
""", unsafe_allow_html=True)

if df_raw is None:
    # Landing state
    col1, col2, col3 = st.columns(3)
    with col1:
        st.markdown("""
        <div class="kpi-card">
            <div class="kpi-value">01</div>
            <div class="kpi-label">Load Your Data</div>
            <div style="font-size:0.85rem;color:#555;margin-top:0.5rem;">
                Upload an Excel workbook or connect to Redshift using the sidebar.
            </div>
        </div>
        """, unsafe_allow_html=True)
    with col2:
        st.markdown("""
        <div class="kpi-card accent">
            <div class="kpi-value magenta">02</div>
            <div class="kpi-label">Filter & Configure</div>
            <div style="font-size:0.85rem;color:#555;margin-top:0.5rem;">
                Filter by LOB, date range, channel, or market. Set your report title and client name.
            </div>
        </div>
        """, unsafe_allow_html=True)
    with col3:
        st.markdown("""
        <div class="kpi-card">
            <div class="kpi-value">03</div>
            <div class="kpi-label">Generate & Download</div>
            <div style="font-size:0.85rem;color:#555;margin-top:0.5rem;">
                One click generates a branded PPTX performance summary deck ready for presentation.
            </div>
        </div>
        """, unsafe_allow_html=True)

    st.markdown("<br>", unsafe_allow_html=True)
    st.info("← Get started by selecting a data source in the sidebar.", icon="👈")

else:
    # ── Column mapping ─────────────────────────────────────────────────────────
    # Normalize common column name variants
    rename_map = {
        "spend": "Spend", "Spend": "Spend",
        "impressions": "Impressions", "Impressions": "Impressions",
        "clicks": "Clicks", "Clicks": "Clicks",
        "site_visits": "Site_Visits", "Site_Visits": "Site_Visits", "visits": "Site_Visits",
        "conversion": "Conversion", "conversions": "Conversion", "Conversion": "Conversion",
        "lob": "LOB", "LOB": "LOB",
        "channel": "Channel", "Channel": "Channel",
        "platform": "Platform", "Platform": "Platform",
        "date": "Date", "Date": "Date",
        "market": "Market", "Market": "Market",
        "product": "Product", "Product": "Product",
        "campaign": "Campaign", "Campaign": "Campaign",
        "business_audience": "Audience", "Business_Audience": "Audience",
        "audience": "Audience", "Audience": "Audience",
        "creative": "Creative", "Creative": "Creative",
        "creative_name": "Creative", "Ad_Name": "Creative", "ad_name": "Creative",
        "ad": "Creative", "Ad": "Creative",
    }
    df = df_raw.rename(columns={k: v for k, v in rename_map.items() if k in df_raw.columns})

    # Ensure numeric cols
    for col in ["Spend", "Impressions", "Clicks", "Site_Visits", "Conversion"]:
        if col in df.columns:
            df[col] = pd.to_numeric(df[col], errors="coerce").fillna(0)

    # ── Filters ────────────────────────────────────────────────────────────────
    st.markdown('<div class="section-header">Filters</div>', unsafe_allow_html=True)
    filter_cols = st.columns(4)

    lob_options = ["All"]
    if "LOB" in df.columns:
        lob_options += sorted(df["LOB"].dropna().unique().tolist())

    with filter_cols[0]:
        selected_lob = st.selectbox("Line of Business", lob_options)

    date_range_label = ""
    if "Date" in df.columns:
        df["Date"] = pd.to_datetime(df["Date"], errors="coerce")
        min_d, max_d = df["Date"].min(), df["Date"].max()
        with filter_cols[1]:
            start_d = st.date_input("Start Date", value=min_d, min_value=min_d, max_value=max_d)
        with filter_cols[2]:
            end_d   = st.date_input("End Date",   value=max_d, min_value=min_d, max_value=max_d)
        date_range_label = f"{start_d.strftime('%b %Y')} – {end_d.strftime('%b %Y')}"

    channel_opts = ["All"]
    if "Channel" in df.columns:
        channel_opts += sorted(df["Channel"].dropna().unique().tolist())
    with filter_cols[3]:
        selected_channel = st.selectbox("Channel", channel_opts)

    # Apply filters
    dff = df.copy()
    if selected_lob != "All" and "LOB" in dff.columns:
        dff = dff[dff["LOB"] == selected_lob]
    if "Date" in dff.columns and "start_d" in dir():
        dff = dff[(dff["Date"] >= pd.Timestamp(start_d)) & (dff["Date"] <= pd.Timestamp(end_d))]
    if selected_channel != "All" and "Channel" in dff.columns:
        dff = dff[dff["Channel"] == selected_channel]

    st.markdown(f"**{len(dff):,} rows** after filters", unsafe_allow_html=False)
    st.markdown("<hr style='border:none;border-top:1px solid #E0E0E0;margin:0.5rem 0 1.5rem 0'>", unsafe_allow_html=True)

    # ── KPI Banner ─────────────────────────────────────────────────────────────
    kpis = compute_kpis(dff)

    k1, k2, k3, k4, k5, k6 = st.columns(6)
    kpi_cards = [
        (k1, fmt_currency(kpis["spend"]),       "Total Spend",      "navy"),
        (k2, fmt_int(kpis["impressions"]),       "Impressions",      "navy"),
        (k3, fmt_int(kpis["clicks"]),            "Clicks",           "navy"),
        (k4, fmt_int(kpis["conversions"]),       "Conversions",      "magenta"),
        (k5, fmt_pct(kpis["ctr"]),               "CTR",              "navy"),
        (k6, fmt_currency(kpis["cpa"],2) if kpis["cpa"] else "—", "CPA", "magenta"),
    ]
    for col, val, label, style in kpi_cards:
        accent = "accent" if style == "magenta" else ""
        val_class = "kpi-value magenta" if style == "magenta" else "kpi-value"
        with col:
            st.markdown(f"""
            <div class="kpi-card {accent}">
                <div class="{val_class}">{val}</div>
                <div class="kpi-label">{label}</div>
            </div>
            """, unsafe_allow_html=True)

    st.markdown("<br>", unsafe_allow_html=True)

    # ── Tabs ───────────────────────────────────────────────────────────────────
    tab_trends, tab1, tab2, tab3, tab5, tab6, tab7 = st.tabs([
        "📈 Trends", "📊 Channel", "🖥 Platform", "📦 Product", "💡 Insights", "🗂 All Data", "🧪 AI Sandbox"
    ])

    with tab1:
        st.markdown('<div class="section-header">Channel Performance</div>', unsafe_allow_html=True)
        if "Channel" in dff.columns:
            sv_col = "Site_Visits" if "Site_Visits" in dff.columns else None
            ch_agg = {"Spend":("Spend","sum"),"Impressions":("Impressions","sum"),
                      "Clicks":("Clicks","sum"),"Conversions":("Conversion","sum")}
            if sv_col: ch_agg["Site_Visits"]=(sv_col,"sum")
            ch_df = dff.groupby("Channel").agg(**ch_agg).reset_index().sort_values("Spend", ascending=False)
            ch_df["CTR"]  = ch_df.apply(lambda r: fmt_pct(r["Clicks"]/r["Impressions"]) if r["Impressions"]>0 else "—", axis=1)
            ch_df["CPA"]  = ch_df.apply(lambda r: fmt_table_currency(r["Spend"]/r["Conversions"]) if r["Conversions"]>0 else "—", axis=1)
            ch_df["CPM"]  = ch_df.apply(lambda r: fmt_table_currency(r["Spend"]/r["Impressions"]*1000) if r["Impressions"]>0 else "—", axis=1)
            ch_df["CPC"]  = ch_df.apply(lambda r: fmt_table_currency(r["Spend"]/r["Clicks"]) if r["Clicks"]>0 else "—", axis=1)
            ch_df["Spend_fmt"] = ch_df["Spend"].apply(fmt_table_currency)
            ch_df["Impressions_fmt"] = ch_df["Impressions"].apply(fmt_table_int)
            ch_df["Conversions_fmt"] = ch_df["Conversions"].apply(fmt_table_int)

            disp_cols = ["Channel","Spend_fmt","Impressions_fmt","Conversions_fmt","CTR","CPA","CPM","CPC"]
            disp_names = ["Channel","Spend","Impressions","Conversions","CTR","CPA","CPM","CPC"]
            if sv_col:
                ch_df["Site_Visits_fmt"] = ch_df["Site_Visits"].apply(fmt_table_int)
                ch_df["CPSV"] = ch_df.apply(lambda r: fmt_table_currency(r["Spend"]/r["Site_Visits"]) if r.get("Site_Visits",0)>0 else "—", axis=1)
                disp_cols += ["Site_Visits_fmt","CPSV"]
                disp_names += ["Site Visits","CPSV"]
            display_ch = ch_df[disp_cols].copy()
            display_ch.columns = disp_names
            st.dataframe(display_ch, use_container_width=True, hide_index=True)

            import plotly.express as px
            fig = px.bar(ch_df, x="Spend", y="Channel", orientation="h",
                         color_discrete_sequence=["#1A2B4A"],
                         title="Spend by Channel")
            fig.update_layout(
                plot_bgcolor="white", paper_bgcolor="white",
                font=dict(family="Avenir, Calibri", color="#333333"),
                title_font_size=14, title_font_color="#1A2B4A",
                xaxis=dict(title="", tickfont=dict(color="#333333", size=12), showgrid=True, gridcolor="#EEEEEE"),
                yaxis=dict(title="", tickfont=dict(color="#333333", size=12)),
                margin=dict(l=0,r=0,t=40,b=0),
                height=300,
            )
            fig.update_traces(marker_line_width=0)
            st.plotly_chart(fig, use_container_width=True)
        else:
            st.info("No 'Channel' column found in this dataset.")

    with tab2:
        st.markdown('<div class="section-header">Platform Performance</div>', unsafe_allow_html=True)
        if "Platform" in dff.columns:
            sv_col2 = "Site_Visits" if "Site_Visits" in dff.columns else None
            pl_agg = {"Spend":("Spend","sum"),"Impressions":("Impressions","sum"),
                      "Clicks":("Clicks","sum"),"Conversions":("Conversion","sum")}
            if sv_col2: pl_agg["Site_Visits"]=(sv_col2,"sum")
            pl_df = dff.groupby("Platform").agg(**pl_agg).reset_index().sort_values("Spend", ascending=False)
            pl_df["CTR"]  = pl_df.apply(lambda r: fmt_pct(r["Clicks"]/r["Impressions"]) if r["Impressions"]>0 else "—", axis=1)
            pl_df["CPA"]  = pl_df.apply(lambda r: fmt_table_currency(r["Spend"]/r["Conversions"]) if r["Conversions"]>0 else "—", axis=1)
            pl_df["CPM"]  = pl_df.apply(lambda r: fmt_table_currency(r["Spend"]/r["Impressions"]*1000) if r["Impressions"]>0 else "—", axis=1)
            pl_df["CPC"]  = pl_df.apply(lambda r: fmt_table_currency(r["Spend"]/r["Clicks"]) if r["Clicks"]>0 else "—", axis=1)

            c1, c2 = st.columns([1.2, 1])
            with c1:
                display_pl = pl_df.copy()
                display_pl["Spend"] = display_pl["Spend"].apply(fmt_table_currency)
                display_pl["Impressions"] = display_pl["Impressions"].apply(fmt_table_int)
                display_pl["Conversions"] = display_pl["Conversions"].apply(fmt_table_int)
                pl_disp_cols = ["Platform","Spend","Impressions","Conversions","CTR","CPA","CPM","CPC"]
                pl_disp_names = ["Platform","Spend","Impressions","Conversions","CTR","CPA","CPM","CPC"]
                if sv_col2:
                    display_pl["Site_Visits"] = display_pl["Site_Visits"].apply(fmt_table_int)
                    display_pl["CPSV"] = pl_df.apply(lambda r: fmt_table_currency(r["Spend"]/r["Site_Visits"]) if r.get("Site_Visits",0)>0 else "—", axis=1)
                    pl_disp_cols += ["Site_Visits","CPSV"]
                    pl_disp_names += ["Site Visits","CPSV"]
                display_pl_out = display_pl[pl_disp_cols].copy()
                display_pl_out.columns = pl_disp_names
                st.dataframe(display_pl_out, use_container_width=True, hide_index=True)
            with c2:
                import plotly.express as px
                pl_plot = pl_df[pl_df["Conversions"] > 0].copy()
                if pl_plot.empty:
                    pl_plot = pl_df.copy()
                fig2 = px.scatter(pl_plot, x="Spend", y="Conversions", size="Spend",
                                  text="Platform", color_discrete_sequence=["#DF552C"],
                                  title="Spend vs Conversions by Platform")
                fig2.update_layout(
                    plot_bgcolor="white", paper_bgcolor="white",
                    font=dict(family="Avenir, Calibri", color="#333333"),
                    title_font_size=13, title_font_color="#1A2B4A",
                    xaxis=dict(tickfont=dict(color="#333333", size=11), title=dict(font=dict(color="#333333"))),
                    yaxis=dict(tickfont=dict(color="#333333", size=11), title=dict(font=dict(color="#333333"))),
                    height=280, margin=dict(l=0,r=0,t=40,b=0),
                )
                fig2.update_traces(textposition="top center", marker_line_width=0)
                st.plotly_chart(fig2, use_container_width=True)
        else:
            st.info("No 'Platform' column found in this dataset.")

    with tab3:
        st.markdown('<div class="section-header">Product Performance</div>', unsafe_allow_html=True)
        prod_col = next((c for c in ["Product","product","LOB","lob","Campaign_Initiative"] if c in dff.columns), None)
        if prod_col:
            sv_col3 = "Site_Visits" if "Site_Visits" in dff.columns else None
            pr_agg = {"Spend":("Spend","sum"),"Conversions":("Conversion","sum"),
                      "Clicks":("Clicks","sum"),"Impressions":("Impressions","sum")}
            if sv_col3: pr_agg["Site_Visits"]=(sv_col3,"sum")
            pr_df = dff.groupby(prod_col).agg(**pr_agg).reset_index().sort_values("Spend", ascending=False).head(10)
            pr_df["CPA"]  = pr_df.apply(lambda r: r["Spend"]/r["Conversions"] if r["Conversions"]>0 else None, axis=1)
            pr_df["CPA_fmt"]   = pr_df["CPA"].apply(lambda v: fmt_table_currency(v) if v else "—")
            pr_df["CPM_fmt"]   = pr_df.apply(lambda r: fmt_table_currency(r["Spend"]/r["Impressions"]*1000) if r["Impressions"]>0 else "—", axis=1)
            pr_df["CPC_fmt"]   = pr_df.apply(lambda r: fmt_table_currency(r["Spend"]/r["Clicks"]) if r["Clicks"]>0 else "—", axis=1)
            pr_df["Spend_fmt"] = pr_df["Spend"].apply(fmt_table_currency)
            pr_df["Conv_fmt"]  = pr_df["Conversions"].apply(fmt_table_int)

            pr_disp = {prod_col: "Segment", "Spend_fmt": "Spend", "Conv_fmt": "Conversions",
                       "CPA_fmt": "CPA", "CPM_fmt": "CPM", "CPC_fmt": "CPC"}
            if sv_col3:
                pr_df["SV_fmt"]   = pr_df["Site_Visits"].apply(fmt_table_int)
                pr_df["CPSV_fmt"] = pr_df.apply(lambda r: fmt_table_currency(r["Spend"]/r["Site_Visits"]) if r.get("Site_Visits",0)>0 else "—", axis=1)
                pr_disp.update({"SV_fmt": "Site Visits", "CPSV_fmt": "CPSV"})

            c1, c2 = st.columns([1, 1.2])
            with c1:
                st.dataframe(
                    pr_df[list(pr_disp.keys())].rename(columns=pr_disp),
                    use_container_width=True, hide_index=True
                )
            with c2:
                import plotly.express as px
                pr_plot = pr_df.dropna(subset=["CPA"]).head(8)
                if not pr_plot.empty:
                    fig3 = px.bar(pr_plot, x=prod_col, y="CPA",
                                  color_discrete_sequence=["#DF552C"],
                                  title="Cost Per Acquisition by Segment")
                    fig3.update_layout(
                        plot_bgcolor="white", paper_bgcolor="white",
                        font=dict(family="Avenir, Calibri", color="#333333"),
                        title_font_size=13, title_font_color="#1A2B4A",
                        xaxis=dict(title="", tickfont=dict(color="#333333", size=11), showgrid=False),
                        yaxis=dict(title=dict(text="CPA ($)", font=dict(color="#333333")), tickfont=dict(color="#333333", size=11), showgrid=True, gridcolor="#EEEEEE"),
                        height=280, margin=dict(l=0,r=0,t=40,b=0),
                    )
                    fig3.update_traces(marker_line_width=0)
                    st.plotly_chart(fig3, use_container_width=True)
        else:
            st.info("No product/segment column detected.")

        # ── Channel > Platform > Creative breakdown ────────────────────────────
        _cpc_cols = [c for c in ["Channel", "Platform", "Creative"] if c in dff.columns]
        if _cpc_cols:
            st.markdown('<div class="section-header" style="margin-top:1.5rem;">Performance by Channel › Platform › Creative</div>', unsafe_allow_html=True)

            _sv_prod = "Site_Visits" if "Site_Visits" in dff.columns else None
            _has_creative = "Creative" in dff.columns

            _cpc_metrics = [
                ("Spend",       "currency"),
                ("Impressions", "int"),
                ("Clicks",      "int"),
                ("CTR",         "pct"),
                ("CPM",         "currency"),
                ("CPA",         "currency"),
            ]
            if _sv_prod:
                _cpc_metrics.append(("Visits", "int"))

            def _fmt_m(val, kind):
                if val is None or (isinstance(val, float) and (val != val)):
                    return "—"
                if kind == "currency": return fmt_table_currency(val)
                if kind == "int":      return fmt_table_int(val)
                if kind == "pct":      return f"{val:.2%}"
                return str(val)

            # Build header
            _cpc_hdrs = [""] + [m for m, _ in _cpc_metrics]
            _cpc_th = "".join(
                f"<th style='background:#1A2B4A;color:#fff;padding:0.45rem 0.6rem;"
                f"font-size:0.7rem;font-weight:700;text-transform:uppercase;"
                f"letter-spacing:0.05em;white-space:nowrap;"
                f"text-align:{'left' if i==0 else 'right'};'>{h}</th>"
                for i, h in enumerate(_cpc_hdrs)
            )
            _cpc_html = f"<table style='width:100%;border-collapse:collapse;font-size:0.83rem;'><thead><tr>{_cpc_th}</tr></thead><tbody>"

            def _cpc_td(val, style=""):
                return f"<td style='padding:0.4rem 0.6rem;border-bottom:1px solid #F0F0F0;{style}'>{val}</td>"

            def _cpc_agg(df_sub):
                sp  = df_sub["Spend"].sum() if "Spend" in df_sub.columns else None
                imp = df_sub["Impressions"].sum() if "Impressions" in df_sub.columns else None
                cl  = df_sub["Clicks"].sum() if "Clicks" in df_sub.columns else None
                cv  = df_sub["Conversion"].sum() if "Conversion" in df_sub.columns else None
                vi  = df_sub[_sv_prod].sum() if _sv_prod and _sv_prod in df_sub.columns else None
                ctr = cl / imp if imp and imp > 0 else None
                cpm = sp / imp * 1000 if imp and imp > 0 else None
                cpa = sp / cv if cv and cv > 0 else None
                return {"Spend": sp, "Impressions": imp, "Clicks": cl, "CTR": ctr, "CPM": cpm, "CPA": cpa, "Visits": vi}

            def _cpc_row(label, agg_d, row_style, label_extra=""):
                cells = _cpc_td(label, label_extra + row_style)
                for m, kind in _cpc_metrics:
                    cells += _cpc_td(_fmt_m(agg_d.get(m), kind), row_style + "text-align:right;")
                return f"<tr>{cells}</tr>"

            _lob_s  = "background:#EDF0F7;font-weight:700;color:#1A2B4A;"
            _ch_s   = "background:#F7F8FA;color:#333;"
            _pl_s   = "background:#fff;color:#555;"
            _cr_s   = "background:#fff;color:#777;"
            _ind1   = "padding-left:18px;"
            _ind2   = "padding-left:34px;"
            _ind3   = "padding-left:50px;"

            ch_vals = sorted(dff["Channel"].unique(), key=lambda x: -dff[dff["Channel"]==x]["Spend"].sum()) if "Channel" in dff.columns else []
            for ch in ch_vals:
                ch_df = dff[dff["Channel"] == ch]
                ch_agg = _cpc_agg(ch_df)
                _cpc_html += _cpc_row(ch, ch_agg, _lob_s)

                pl_vals = sorted(ch_df["Platform"].unique(), key=lambda x: -ch_df[ch_df["Platform"]==x]["Spend"].sum()) if "Platform" in ch_df.columns else []
                for pl in pl_vals:
                    pl_df = ch_df[ch_df["Platform"] == pl]
                    pl_agg = _cpc_agg(pl_df)
                    _cpc_html += _cpc_row(f"↳ {pl}", pl_agg, _ch_s, _ind1)

                    if _has_creative:
                        cr_vals = sorted(pl_df["Creative"].dropna().unique(), key=lambda x: -pl_df[pl_df["Creative"]==x]["Spend"].sum())
                        for cr in cr_vals[:10]:  # cap at 10 creatives per platform
                            cr_df = pl_df[pl_df["Creative"] == cr]
                            cr_agg = _cpc_agg(cr_df)
                            _cr_label = str(cr)[:60] + ("…" if len(str(cr)) > 60 else "")
                            _cpc_html += _cpc_row(f"↳ {_cr_label}", cr_agg, _cr_s, _ind2)

            _cpc_html += "</tbody></table>"
            st.markdown(_cpc_html, unsafe_allow_html=True)

    with tab_trends:
        _sv = "Site_Visits" if "Site_Visits" in dff.columns else None

        # ── Shared hierarchy helpers ───────────────────────────────────────────
        _hier_cols = [c for c in ["LOB", "Channel", "Platform"] if c in dff.columns]

        def _agg_metrics(df_src, group_cols):
            """Aggregate Spend/Impressions/Clicks/CTR/CPM/Visits by group_cols."""
            if not group_cols or df_src.empty:
                return pd.DataFrame()
            agg = {
                "Spend":       ("Spend",       "sum"),
                "Impressions": ("Impressions", "sum"),
                "Clicks":      ("Clicks",      "sum"),
            }
            if _sv:
                agg["Visits"] = (_sv, "sum")
            g = df_src.groupby(group_cols, sort=False).agg(**agg).reset_index()
            g["CTR"] = g.apply(lambda r: r["Clicks"] / r["Impressions"] if r["Impressions"] > 0 else None, axis=1)
            g["CPM"] = g.apply(lambda r: r["Spend"] / r["Impressions"] * 1000 if r["Impressions"] > 0 else None, axis=1)
            return g

        def _fmt_metric(val, kind):
            if val is None or (isinstance(val, float) and pd.isna(val)):
                return "—"
            if kind == "currency":  return fmt_table_currency(val)
            if kind == "int":       return fmt_table_int(val)
            if kind == "pct":       return f"{val*100:.2f}%"
            return str(val)

        def _pct_chg_html(cur_val, prev_val):
            if prev_val is None or prev_val == 0 or pd.isna(prev_val):
                return "—"
            if cur_val is None or pd.isna(cur_val):
                return "—"
            chg = (cur_val - prev_val) / abs(prev_val)
            sign = "+" if chg > 0 else ""
            color = "#2E7D32" if chg > 0 else "#C62828"
            return f"<span style='color:{color};font-weight:600'>{sign}{chg*100:.1f}%</span>"

        # Row style helpers
        _lob_style  = "background:#EDF0F7;font-weight:700;color:#1A2B4A;"
        _ch_style   = "background:#F7F8FA;color:#333;"
        _pl_style   = "background:#fff;color:#555;"
        _indent1    = "padding-left:20px;"
        _indent2    = "padding-left:38px;"

        def _build_hier_html(df_cur, df_prev=None):
            """
            Build an HTML table with LOB > Channel > Platform grouping.
            df_prev: if provided, add % change columns after each metric.
            """
            has_prev  = df_prev is not None and not df_prev.empty
            has_visits = _sv is not None

            metrics = [
                ("Spend",       "currency"),
                ("Impressions", "int"),
                ("Clicks",      "int"),
                ("CTR",         "pct"),
                ("CPM",         "currency"),
            ]
            if has_visits:
                metrics.append(("Visits", "int"))

            # Build header
            hdrs = [""]
            for m, _ in metrics:
                hdrs.append(m)
                if has_prev:
                    hdrs.append("Δ")

            th_inner = "".join(
                f"<th style='background:#1A2B4A;color:#fff;padding:0.45rem 0.6rem;"
                f"font-size:0.7rem;font-weight:700;text-transform:uppercase;"
                f"letter-spacing:0.05em;white-space:nowrap;"
                f"text-align:{'left' if i == 0 else 'right'};'>{h}</th>"
                for i, h in enumerate(hdrs)
            )
            html = f"<table style='width:100%;border-collapse:collapse;font-size:0.83rem;'><thead><tr>{th_inner}</tr></thead><tbody>"

            def td(val, style=""):
                return f"<td style='padding:0.45rem 0.6rem;border-bottom:1px solid #F0F0F0;{style}'>{val}</td>"

            def make_row(label, label_style, row_style, g_cur, g_prev_row):
                cells = td(label, label_style + row_style)
                for m, kind in metrics:
                    cv = g_cur.get(m) if isinstance(g_cur, dict) else (g_cur[m] if m in g_cur.index else None)
                    cells += td(_fmt_metric(cv, kind), row_style + "text-align:right;")
                    if has_prev:
                        pv = g_prev_row.get(m) if isinstance(g_prev_row, dict) else (g_prev_row[m] if g_prev_row is not None and m in g_prev_row.index else None)
                        cells += td(_pct_chg_html(cv, pv), row_style + "text-align:right;")
                return f"<tr>{cells}</tr>"

            def lookup(df, keys):
                """Find a row in df matching the key dict. Returns series or None."""
                if df is None or df.empty:
                    return None
                mask = pd.Series([True] * len(df), index=df.index)
                for col, val in keys.items():
                    if col in df.columns:
                        mask &= df[col] == val
                result = df[mask]
                return result.iloc[0] if not result.empty else None

            # Determine available hierarchy columns in this data
            h_cols = [c for c in ["LOB", "Channel", "Platform"] if c in df_cur.columns]

            if len(h_cols) == 0:
                return "<p style='color:#888;'>No grouping columns available.</p>"

            lob_vals = df_cur["LOB"].unique() if "LOB" in df_cur.columns else [None]

            for lob in sorted(lob_vals, key=lambda x: (x is None, x)):
                lob_filter = {"LOB": lob} if "LOB" in df_cur.columns and lob is not None else {}
                lob_cur  = df_cur[df_cur["LOB"] == lob] if lob_filter else df_cur
                lob_prev = df_prev[df_prev["LOB"] == lob] if (has_prev and "LOB" in df_prev.columns and lob is not None) else df_prev

                # LOB totals
                lob_cur_tot  = {m: lob_cur[m].sum()  if m in lob_cur.columns  else None for m, _ in metrics if m not in ("CTR","CPM")}
                lob_prev_tot = {m: lob_prev[m].sum() if lob_prev is not None and m in lob_prev.columns else None for m, _ in metrics if m not in ("CTR","CPM")} if has_prev else {}
                # Recalc rates for LOB total
                lob_cur_tot["CTR"] = lob_cur["Clicks"].sum() / lob_cur["Impressions"].sum() if "Impressions" in lob_cur.columns and lob_cur["Impressions"].sum() > 0 else None
                lob_cur_tot["CPM"] = lob_cur["Spend"].sum() / lob_cur["Impressions"].sum() * 1000 if "Impressions" in lob_cur.columns and lob_cur["Impressions"].sum() > 0 else None
                if has_prev and lob_prev is not None and not lob_prev.empty:
                    lob_prev_tot["CTR"] = lob_prev["Clicks"].sum() / lob_prev["Impressions"].sum() if lob_prev["Impressions"].sum() > 0 else None
                    lob_prev_tot["CPM"] = lob_prev["Spend"].sum() / lob_prev["Impressions"].sum() * 1000 if lob_prev["Impressions"].sum() > 0 else None

                lob_label = str(lob) if lob is not None else "(All)"
                html += make_row(lob_label, _lob_style, "", lob_cur_tot, lob_prev_tot if has_prev else None)

                if "Channel" not in df_cur.columns:
                    continue

                ch_vals = lob_cur["Channel"].unique()
                for ch in sorted(ch_vals, key=lambda x: -lob_cur[lob_cur["Channel"]==x]["Spend"].sum()):
                    ch_cur  = lob_cur[lob_cur["Channel"] == ch]
                    ch_prev = lob_prev[lob_prev["Channel"] == ch] if (has_prev and lob_prev is not None and "Channel" in lob_prev.columns) else None

                    ch_cur_tot = {m: ch_cur[m].sum() if m in ch_cur.columns else None for m, _ in metrics if m not in ("CTR","CPM")}
                    ch_prev_tot = {m: ch_prev[m].sum() if ch_prev is not None and m in ch_prev.columns else None for m, _ in metrics if m not in ("CTR","CPM")} if has_prev else {}
                    ch_cur_tot["CTR"] = ch_cur["Clicks"].sum() / ch_cur["Impressions"].sum() if "Impressions" in ch_cur.columns and ch_cur["Impressions"].sum() > 0 else None
                    ch_cur_tot["CPM"] = ch_cur["Spend"].sum() / ch_cur["Impressions"].sum() * 1000 if "Impressions" in ch_cur.columns and ch_cur["Impressions"].sum() > 0 else None
                    if has_prev and ch_prev is not None and not ch_prev.empty:
                        ch_prev_tot["CTR"] = ch_prev["Clicks"].sum() / ch_prev["Impressions"].sum() if ch_prev["Impressions"].sum() > 0 else None
                        ch_prev_tot["CPM"] = ch_prev["Spend"].sum() / ch_prev["Impressions"].sum() * 1000 if ch_prev["Impressions"].sum() > 0 else None

                    html += make_row(f"↳ {ch}", _indent1, _ch_style, ch_cur_tot, ch_prev_tot if has_prev else None)

                    if "Platform" not in df_cur.columns:
                        continue

                    pl_vals = ch_cur["Platform"].unique()
                    for pl in sorted(pl_vals, key=lambda x: -ch_cur[ch_cur["Platform"]==x]["Spend"].sum()):
                        pl_cur_row  = ch_cur[ch_cur["Platform"] == pl]
                        pl_prev_row = ch_prev[ch_prev["Platform"] == pl] if (has_prev and ch_prev is not None and "Platform" in ch_prev.columns) else None

                        pl_cur_tot = {m: pl_cur_row[m].sum() if m in pl_cur_row.columns else None for m, _ in metrics if m not in ("CTR","CPM")}
                        pl_prev_tot = {m: pl_prev_row[m].sum() if pl_prev_row is not None and m in pl_prev_row.columns else None for m, _ in metrics if m not in ("CTR","CPM")} if has_prev else {}
                        pl_cur_tot["CTR"] = pl_cur_row["Clicks"].sum() / pl_cur_row["Impressions"].sum() if "Impressions" in pl_cur_row.columns and pl_cur_row["Impressions"].sum() > 0 else None
                        pl_cur_tot["CPM"] = pl_cur_row["Spend"].sum() / pl_cur_row["Impressions"].sum() * 1000 if "Impressions" in pl_cur_row.columns and pl_cur_row["Impressions"].sum() > 0 else None
                        if has_prev and pl_prev_row is not None and not pl_prev_row.empty:
                            pl_prev_tot["CTR"] = pl_prev_row["Clicks"].sum() / pl_prev_row["Impressions"].sum() if pl_prev_row["Impressions"].sum() > 0 else None
                            pl_prev_tot["CPM"] = pl_prev_row["Spend"].sum() / pl_prev_row["Impressions"].sum() * 1000 if pl_prev_row["Impressions"].sum() > 0 else None

                        html += make_row(f"↳ {pl}", _indent2, _pl_style, pl_cur_tot, pl_prev_tot if has_prev else None)

            html += "</tbody></table>"
            return html

        # ── Monthly Trends ─────────────────────────────────────────────────────
        st.markdown('<div class="section-header">Monthly Trends</div>', unsafe_allow_html=True)
        if "Date" in dff.columns:
            monthly = dff.copy()
            monthly["Month"] = monthly["Date"].dt.to_period("M")

            trend = monthly.groupby("Month").agg(
                Spend=("Spend","sum"), Conversions=("Conversion","sum"),
                Clicks=("Clicks","sum"), Impressions=("Impressions","sum"),
                Days=("Date", lambda x: x.dt.date.nunique()),
                **( {"Visits": (_sv, "sum")} if _sv else {} ),
            ).reset_index()
            trend["Month_str"] = trend["Month"].astype(str)
            trend["CTR"] = trend.apply(
                lambda r: r["Clicks"]/r["Impressions"] if r["Impressions"]>0 else None, axis=1)
            trend["CPM"] = trend.apply(
                lambda r: r["Spend"]/r["Impressions"]*1000 if r["Impressions"]>0 else None, axis=1)
            trend["CPA"] = trend.apply(
                lambda r: r["Spend"]/r["Conversions"] if r["Conversions"]>0 else None, axis=1)

            # Daily-normalized columns for fair cross-month comparison
            for _mc in ["Spend", "Impressions", "Clicks", "Conversions"] + (["Visits"] if _sv else []):
                trend[f"{_mc}_pd"] = trend.apply(
                    lambda r: r[_mc] / r["Days"] if r["Days"] > 0 else None, axis=1)

            # ── Trend narrative (above chart) ───────────────────────────────────
            if len(trend) >= 2:
                _t  = trend.sort_values("Month")
                _L  = _t.iloc[-1]   # most recent month
                _P  = _t.iloc[-2]   # prior month
                _F  = _t.iloc[0]    # first month

                # Determine if most recent month is partial (fewer days than prior)
                _days_ratio = _L["Days"] / _P["Days"] if _P["Days"] > 0 else 1.0
                _is_partial = _days_ratio < 0.85  # less than 85% of prior month's data days

                # Use daily-normalized for MoM when partial, raw otherwise
                def _mom(cur_col, prev_val):
                    c = _L[cur_col + ("_pd" if _is_partial else "")]
                    p = _P[prev_val + ("_pd" if _is_partial else "")]
                    return (c - p) / p if p and p > 0 else None

                _mom_spend = _mom("Spend", "Spend")
                _mom_conv  = _mom("Conversions", "Conversions")
                _mom_ctr   = ((_L["CTR"] or 0) - (_P["CTR"] or 0)) / _P["CTR"] if _P["CTR"] and _P["CTR"] > 0 else None
                _mom_cpm   = ((_L["CPM"] or 0) - (_P["CPM"] or 0)) / _P["CPM"] if _P["CPM"] and _P["CPM"] > 0 else None

                _overall_chg = (_L["Spend"] - _F["Spend"]) / _F["Spend"] if _F["Spend"] > 0 else None
                _total_sp    = trend["Spend"].sum()
                _total_conv  = trend["Conversions"].sum()

                peak_spend_row = trend.loc[trend["Spend"].idxmax()]
                peak_conv_row  = trend.loc[trend["Conversions"].idxmax()]

                def _arrow(v):
                    if v is None: return ""
                    return "▲" if v > 0 else "▼"
                def _chg(v, invert=False):
                    if v is None: return "—"
                    good = (v > 0) != invert
                    color = "#1B7F4C" if good else "#C0392B"
                    return f"<span style='color:{color};font-weight:700;'>{_arrow(v)} {abs(v):.0%}</span>"

                # Build narrative — most recent month first
                _L_days, _P_days = int(_L["Days"]), int(_P["Days"])
                _partial_note = f" (partial month — {_L_days} days vs. {_P_days} prior; rates normalized per day)" if _is_partial else ""
                _basis = "daily-rate basis" if _is_partial else "vs. prior month"

                _bullets = []

                # 1. Most recent month headline
                _sp_chg_str = _chg(_mom_spend)
                _cv_chg_str = _chg(_mom_conv)
                _ctr_str    = _chg(_mom_ctr)
                _cpm_str    = _chg(_mom_cpm, invert=True)

                _bullets.append(
                    f"<b>{_L['Month_str']}{_partial_note}</b> — "
                    f"Spend {fmt_currency(_L['Spend'])} ({_sp_chg_str} {_basis}); "
                    f"{fmt_int(int(_L['Conversions']))} conversions ({_cv_chg_str}); "
                    f"CTR {_L['CTR']:.2%} ({_ctr_str}); "
                    f"CPM {fmt_currency(_L['CPM'])} ({_cpm_str})."
                    + (f" CPA {fmt_currency(_L['CPA'])}." if _L.get('CPA') else "")
                )

                # 2. Period context
                if _overall_chg is not None:
                    _dir = "up" if _overall_chg > 0 else "down"
                    _bullets.append(
                        f"Spend trended {_dir} {abs(_overall_chg):.0%} from {_F['Month_str']} to {_L['Month_str']}, "
                        f"totalling {fmt_currency(_total_sp)} ({fmt_int(int(_total_conv))} conversions) across {len(_t)} months."
                    )

                # 3. Peak months
                _bullets.append(
                    f"Highest spend: {peak_spend_row['Month_str']} ({fmt_currency(peak_spend_row['Spend'])}). "
                    f"Highest conversions: {peak_conv_row['Month_str']} ({fmt_int(int(peak_conv_row['Conversions']))})."
                )

                _html_bullets = "".join(f"<li style='margin-bottom:0.35rem;'>{b}</li>" for b in _bullets)
                st.markdown(
                    f"<ul style='color:#444;font-size:0.87rem;line-height:1.65;margin:0.3rem 0 1.2rem 1rem;"
                    f"padding:0;list-style:disc;'>{_html_bullets}</ul>",
                    unsafe_allow_html=True,
                )
            else:
                peak_spend_row = trend.loc[trend["Spend"].idxmax()]
                peak_conv_row  = trend.loc[trend["Conversions"].idxmax()]

            # Selector for secondary axis — right-aligned
            _metric_options = ["Conversions", "Impressions", "CTR (%)", "CPM ($)"]
            if _sv:
                _metric_options.append("Visits")
            _, _sel_col = st.columns([5, 2])
            with _sel_col:
                _secondary = st.selectbox(
                    "Secondary Metric",
                    _metric_options,
                    index=0,
                    key="trend_secondary",
                )

            # Map selection to column + formatting
            _metric_map = {
                "Conversions": ("Conversions", "Conversions",  "#DF552C", "{:,.0f}"),
                "Impressions": ("Impressions", "Impressions",  "#2E6DA4", "{:,.0f}"),
                "CTR (%)":     ("CTR",         "CTR (%)",      "#2E8B57", "{:.2%}"),
                "CPM ($)":     ("CPM",         "CPM ($)",      "#8B6914", "${:,.2f}"),
                "Visits":      ("Visits",       "Visits",       "#6B4CA0", "{:,.0f}"),
            }
            _col, _label, _color, _fmt = _metric_map[_secondary]

            import plotly.graph_objects as go
            fig4 = go.Figure()
            fig4.add_trace(go.Bar(
                x=trend["Month_str"], y=trend["Spend"],
                name="Spend", marker_color="#1A2B4A", yaxis="y1",
            ))
            fig4.add_trace(go.Scatter(
                x=trend["Month_str"], y=trend[_col],
                name=_label, line=dict(color=_color, width=3),
                mode="lines+markers", marker=dict(size=8), yaxis="y2",
            ))

            # Tick formatting for right axis
            _y2_tickfmt = ".2%" if _col == "CTR" else ("$,.2f" if _col == "CPM" else ",.0f")
            _y2_tickpfx = "$" if _col == "CPM" else ""

            fig4.update_layout(
                plot_bgcolor="white", paper_bgcolor="white",
                font=dict(family="Avenir, Calibri", color="#333333"),
                title=f"Monthly Spend vs. {_label}",
                title_font_size=14, title_font_color="#1A2B4A",
                xaxis=dict(tickfont=dict(color="#333333", size=11), showgrid=False),
                yaxis=dict(
                    title=dict(text="Spend ($)", font=dict(color="#333333")),
                    tickfont=dict(color="#333333", size=11),
                    showgrid=True, gridcolor="#EEEEEE",
                ),
                yaxis2=dict(
                    title=dict(text=_label, font=dict(color=_color)),
                    tickfont=dict(color=_color, size=11),
                    tickformat=_y2_tickfmt,
                    tickprefix=_y2_tickpfx,
                    overlaying="y", side="right",
                ),
                legend=dict(orientation="h", yanchor="bottom", y=1.02, font=dict(color="#333333")),
                margin=dict(l=0, r=0, t=60, b=0),
                height=380,
            )
            st.plotly_chart(fig4, use_container_width=True)

            m1, m2 = st.columns(2)
            with m1:
                st.markdown(f"""
                <div class="insight-card">
                    <div class="insight-title">Peak Spend Month</div>
                    <div class="insight-text">{peak_spend_row["Month_str"]} — {fmt_currency(peak_spend_row["Spend"])}</div>
                </div>
                """, unsafe_allow_html=True)
            with m2:
                st.markdown(f"""
                <div class="insight-card">
                    <div class="insight-title">Peak Conversion Month</div>
                    <div class="insight-text">{peak_conv_row["Month_str"]} — {fmt_int(peak_conv_row["Conversions"])} conversions</div>
                </div>
                """, unsafe_allow_html=True)

            # ── Stacked bar: Spend by Channel over time ────────────────────────
            if "Channel" in dff.columns:
                _ch_hdr, _, _ch_sel = st.columns([4, 1, 2])
                with _ch_hdr:
                    st.markdown("<p style='font-weight:700;color:#1A2B4A;margin-top:1.8rem;margin-bottom:0.4rem;font-size:0.95rem;'>Monthly Spend by Channel</p>", unsafe_allow_html=True)
                with _ch_sel:
                    _ch_sec_opts = ["Conversions", "Impressions", "CTR (%)", "CPM ($)"]
                    if _sv:
                        _ch_sec_opts.append("Visits")
                    _ch_secondary = st.selectbox(
                        "Secondary Metric",
                        _ch_sec_opts,
                        index=0,
                        key="ch_trend_secondary",
                    )

                ch_trend = dff.copy()
                ch_trend["Month"] = ch_trend["Date"].dt.to_period("M")

                # Aggregate all metrics by Month + Channel
                _ch_agg_spec = {
                    "Spend":       ("Spend",      "sum"),
                    "Impressions": ("Impressions", "sum"),
                    "Clicks":      ("Clicks",      "sum"),
                    "Conversions": ("Conversion",  "sum"),
                }
                if _sv:
                    _ch_agg_spec["Visits"] = (_sv, "sum")
                ch_monthly = (
                    ch_trend.groupby(["Month", "Channel"])
                    .agg(**_ch_agg_spec)
                    .reset_index()
                )
                ch_monthly["Month_str"] = ch_monthly["Month"].astype(str)
                ch_monthly = ch_monthly.sort_values("Month_str")

                # Monthly totals for secondary line
                _mo_totals = ch_monthly.groupby("Month_str").agg(
                    Spend=("Spend","sum"), Impressions=("Impressions","sum"),
                    Clicks=("Clicks","sum"), Conversions=("Conversions","sum"),
                    **( {"Visits": ("Visits","sum")} if _sv else {} ),
                ).reset_index()
                _mo_totals["CTR"] = _mo_totals.apply(
                    lambda r: r["Clicks"]/r["Impressions"] if r["Impressions"]>0 else None, axis=1)
                _mo_totals["CPM"] = _mo_totals.apply(
                    lambda r: r["Spend"]/r["Impressions"]*1000 if r["Impressions"]>0 else None, axis=1)

                # Map secondary selection
                _ch_metric_map = {
                    "Conversions": ("Conversions", "Conversions", "#DF552C", ",.0f",  ""),
                    "Impressions": ("Impressions", "Impressions", "#2E6DA4", ",.0f",  ""),
                    "CTR (%)":     ("CTR",         "CTR (%)",     "#2E8B57", ".2%",   ""),
                    "CPM ($)":     ("CPM",         "CPM ($)",     "#8B6914", ",.2f",  "$"),
                    "Visits":      ("Visits",       "Visits",      "#6B4CA0", ",.0f",  ""),
                }
                _ch_col, _ch_lbl, _ch_clr, _ch_fmt, _ch_pfx = _ch_metric_map[_ch_secondary]

                # Order channels by total spend for consistent stacking
                ch_order = (
                    ch_monthly.groupby("Channel")["Spend"]
                    .sum().sort_values(ascending=False).index.tolist()
                )
                _ch_palette = [
                    "#1A2B4A", "#DF552C", "#2E6DA4", "#E8954A",
                    "#4A90B8", "#8C3D1E", "#6B8FC4", "#A0522D",
                    "#264E73", "#C4753A",
                ]
                ch_colors = {ch: _ch_palette[i % len(_ch_palette)] for i, ch in enumerate(ch_order)}

                fig_ch_trend = go.Figure()
                # Stacked spend bars
                for ch in ch_order:
                    ch_data = ch_monthly[ch_monthly["Channel"] == ch]
                    fig_ch_trend.add_trace(go.Bar(
                        x=ch_data["Month_str"],
                        y=ch_data["Spend"],
                        name=ch,
                        marker_color=ch_colors[ch],
                        yaxis="y1",
                        hovertemplate=f"<b>{ch}</b><br>%{{x}}<br>Spend: $%{{y:,.0f}}<extra></extra>",
                    ))
                # Secondary metric line (total across all channels)
                fig_ch_trend.add_trace(go.Scatter(
                    x=_mo_totals["Month_str"],
                    y=_mo_totals[_ch_col],
                    name=_ch_lbl,
                    line=dict(color=_ch_clr, width=3, dash="solid"),
                    mode="lines+markers", marker=dict(size=7),
                    yaxis="y2",
                    hovertemplate=f"<b>{_ch_lbl}</b><br>%{{x}}<br>%{{y:{_ch_fmt}}}<extra></extra>",
                ))

                fig_ch_trend.update_layout(
                    barmode="stack",
                    plot_bgcolor="white",
                    paper_bgcolor="white",
                    font=dict(family="Avenir, Calibri", color="#333333"),
                    xaxis=dict(tickfont=dict(color="#333333", size=11), showgrid=False),
                    yaxis=dict(
                        title=dict(text="Spend ($)", font=dict(color="#333333")),
                        tickfont=dict(color="#333333", size=11),
                        showgrid=True, gridcolor="#EEEEEE",
                        tickprefix="$", tickformat=",.0f",
                    ),
                    yaxis2=dict(
                        title=dict(text=_ch_lbl, font=dict(color=_ch_clr)),
                        tickfont=dict(color=_ch_clr, size=11),
                        tickformat=_ch_fmt,
                        tickprefix=_ch_pfx,
                        overlaying="y", side="right", showgrid=False,
                    ),
                    legend=dict(
                        orientation="h", yanchor="bottom", y=1.02,
                        font=dict(color="#333333", size=11),
                    ),
                    margin=dict(l=10, r=10, t=50, b=10),
                    height=400,
                )
                st.plotly_chart(fig_ch_trend, use_container_width=True)

            # Hierarchical monthly summary table
            if _hier_cols:
                st.markdown("<p style='font-weight:700;color:#1A2B4A;margin-top:1.5rem;margin-bottom:0.4rem;font-size:0.95rem;'>Monthly Summary — LOB › Channel › Platform</p>", unsafe_allow_html=True)
                mo_hier = _agg_metrics(dff, _hier_cols)
                if not mo_hier.empty:
                    st.markdown(_build_hier_html(mo_hier), unsafe_allow_html=True)
        else:
            st.info("No 'Date' column found for trend analysis.")

        st.markdown("<hr style='margin:2rem 0;border-color:#E0E0E0;'>", unsafe_allow_html=True)

        # ── Weekly Performance Summary ─────────────────────────────────────────
        st.markdown('<div class="section-header">Weekly Performance Summary</div>', unsafe_allow_html=True)
        if "Date" in dff.columns and len(dff) > 0:
            latest_date = dff["Date"].max()
            wd = latest_date.weekday()
            days_back   = (wd + 1) % 7
            week_end    = latest_date - pd.Timedelta(days=days_back)
            week_start  = week_end - pd.Timedelta(days=6)
            prior_end   = week_start - pd.Timedelta(days=1)
            prior_start = prior_end - pd.Timedelta(days=6)

            cur_wk  = dff[(dff["Date"] >= week_start)  & (dff["Date"] <= week_end)]
            prev_wk = dff[(dff["Date"] >= prior_start) & (dff["Date"] <= prior_end)]

            st.markdown(
                f"<p style='color:#555;font-size:0.85rem;margin:-0.5rem 0 1rem 0;'>"
                f"Most recent complete week: <b>{week_start.strftime('%b %-d')} – {week_end.strftime('%b %-d, %Y')}</b> "
                f"vs. prior week: {prior_start.strftime('%b %-d')} – {prior_end.strftime('%b %-d')}"
                f"</p>",
                unsafe_allow_html=True,
            )

            if _hier_cols and not cur_wk.empty:
                wk_cur_hier  = _agg_metrics(cur_wk,  _hier_cols)
                wk_prev_hier = _agg_metrics(prev_wk, _hier_cols) if not prev_wk.empty else pd.DataFrame()
                st.markdown(_build_hier_html(wk_cur_hier, wk_prev_hier if not wk_prev_hier.empty else None), unsafe_allow_html=True)
            else:
                st.info("No data for the most recent complete week.")

    with tab5:
        st.markdown('<div class="section-header">Automated Insights</div>', unsafe_allow_html=True)
        insights = generate_insights(dff, kpis)
        if insights:
            for ins in insights:
                st.markdown(f"""
                <div class="insight-card">
                    <div class="insight-title">{ins["title"]}</div>
                    <div class="insight-text">{ins["text"]}</div>
                </div>
                """, unsafe_allow_html=True)
        else:
            st.info("Load data to generate insights.")

    # ── Tab 6: All Data ────────────────────────────────────────────────────────
    with tab6:
        st.markdown('<div class="section-header">All Data</div>', unsafe_allow_html=True)

        search_col, dl_col = st.columns([3, 1])
        with search_col:
            search_term = st.text_input("Search", placeholder="Filter rows by keyword…", label_visibility="collapsed")
        with dl_col:
            csv_bytes = dff.to_csv(index=False).encode()
            st.download_button("⬇ Export CSV", data=csv_bytes,
                               file_name="campaign_data.csv", mime="text/csv",
                               use_container_width=True)

        display_df = dff.copy()
        if search_term:
            mask = display_df.apply(
                lambda col: col.astype(str).str.contains(search_term, case=False, na=False)
            ).any(axis=1)
            display_df = display_df[mask]

        st.markdown(f"<span style='font-size:0.8rem;color:#8C8C8C;'>{len(display_df):,} rows</span>", unsafe_allow_html=True)
        st.dataframe(display_df, use_container_width=True, height=520)

    # ── Tab 7: AI Sandbox ──────────────────────────────────────────────────────
    with tab7:
        st.markdown('<div class="section-header">AI Visualization Sandbox</div>', unsafe_allow_html=True)
        st.markdown(
            "<p style='color:#555;font-size:0.9rem;margin-bottom:1.25rem;'>"
            "Ask a question about the data in plain English. The AI will generate a chart and summary table."
            "</p>",
            unsafe_allow_html=True,
        )

        key_col, _ = st.columns([2, 3])
        with key_col:
            anthropic_key = st.text_input(
                "Anthropic API Key",
                type="password",
                placeholder="sk-ant-...",
                help="Get yours at console.anthropic.com",
                key="api_key_input",
            )
            if anthropic_key:
                st.session_state["anthropic_api_key"] = anthropic_key
                _save_config({"anthropic_api_key": anthropic_key})

        # Example prompts
        examples = [
            "Show spend and conversions by platform",
            "What is the CPA by product, sorted best to worst?",
            "Compare CTR for brand vs non-brand campaigns",
            "Show monthly conversion trend with spend overlay",
            "Which channels have the highest CVR?",
        ]
        st.markdown("<div style='margin-bottom:0.75rem;'>", unsafe_allow_html=True)
        ex_cols = st.columns(len(examples))
        for i, ex in enumerate(examples):
            with ex_cols[i]:
                if st.button(ex, key=f"ex_{i}", use_container_width=True):
                    st.session_state["sandbox_query"] = ex
        st.markdown("</div>", unsafe_allow_html=True)

        query = st.text_input(
            "Your question",
            value=st.session_state.get("sandbox_query", ""),
            placeholder="e.g. Show me spend and CPA by channel",
            label_visibility="collapsed",
            key="sandbox_input",
        )

        run_query = st.button("▶ Run", use_container_width=False)

        if run_query and query:
            # Build a compact data summary for the prompt
            col_list = ", ".join(dff.columns.tolist())
            sample = dff.head(5).to_csv(index=False)
            numeric_cols = dff.select_dtypes(include="number").columns.tolist()
            cat_cols = [c for c in dff.columns if c not in numeric_cols]

            system_prompt = """You are a data analyst writing Python code for a Streamlit app.
You will be given a pandas DataFrame called `dff` and a plain-English question.
Return ONLY valid Python code — no markdown, no explanation, no code fences.

The code must:
1. Compute an aggregation or analysis from `dff` to answer the question.
   Store the result in a variable called `result_df` (a pandas DataFrame).
2. Create a Plotly figure stored in a variable called `fig`.
   Use these exact colors: primary #1A2B4A, accent #DF552C, light gray #F4F4F4.
   Set plot_bgcolor='white', paper_bgcolor='white'.
   Set font=dict(family='Avenir, Calibri', color='#333333').
   Set title_font_color='#1A2B4A'.
   Keep the chart clean and readable. Height 380.
3. Do NOT call st.plotly_chart() or st.dataframe() — just define fig and result_df.
4. Do NOT import anything — pandas as pd, plotly.express as px, and plotly.graph_objects as go are already available.
5. If the question cannot be answered from the data, set fig=None and result_df=None."""

            user_prompt = f"""DataFrame columns: {col_list}
Numeric columns: {', '.join(numeric_cols)}
Categorical columns: {', '.join(cat_cols)}
Sample rows:
{sample}

Question: {query}"""

            _progress = st.progress(0, text="Checking credentials…")
            try:
                    import anthropic, os
                    api_key = st.session_state.get("anthropic_api_key") or os.environ.get("ANTHROPIC_API_KEY")
                    if not api_key:
                        _progress.empty()
                        st.error("Enter your Anthropic API key above to use the AI Sandbox.")
                        st.stop()

                    _progress.progress(15, text="Sending query to AI…")
                    client_ai = anthropic.Anthropic(api_key=api_key)
                    response = client_ai.messages.create(
                        model="claude-sonnet-4-6",
                        max_tokens=4096,
                        system=system_prompt,
                        messages=[{"role": "user", "content": user_prompt}],
                    )
                    _progress.progress(70, text="Generating visualization…")
                    code = response.content[0].text.strip()
                    # Strip accidental fences
                    if code.startswith("```"):
                        code = "\n".join(code.split("\n")[1:])
                    if code.endswith("```"):
                        code = "\n".join(code.split("\n")[:-1])

                    _progress.progress(85, text="Rendering results…")
                    import plotly.express as px
                    import plotly.graph_objects as go
                    local_ns = {"dff": dff.copy(), "pd": pd, "px": px, "go": go}
                    exec(code, local_ns)
                    fig_out    = local_ns.get("fig")
                    result_out = local_ns.get("result_df")
                    _progress.progress(100, text="Done")
                    _progress.empty()

                    if fig_out is not None:
                        st.plotly_chart(fig_out, use_container_width=True)
                    else:
                        st.warning("No chart was generated for this query.")

                    if result_out is not None and not result_out.empty:
                        st.markdown('<div class="section-header" style="margin-top:1rem;">Summary Table</div>', unsafe_allow_html=True)
                        st.dataframe(result_out, use_container_width=True, hide_index=True)
                    elif result_out is not None and result_out.empty:
                        st.info("Query returned no rows.")

            except ImportError:
                    _progress.empty()
                    st.error("Anthropic SDK not installed. Run: pip install anthropic")
            except Exception as e:
                    _progress.empty()
                    st.error(f"Error running query: {e}")
                    with st.expander("Generated code"):
                        st.code(code if 'code' in locals() else "No code generated", language="python")

    # ── Generate deck ──────────────────────────────────────────────────────────
    st.markdown("<br>", unsafe_allow_html=True)
    st.markdown('<div class="section-header">Generate Deck</div>', unsafe_allow_html=True)

    gen_col1, gen_col2 = st.columns([2, 1])
    with gen_col1:
        lob_label = selected_lob if selected_lob != "All" else report_title
        st.markdown(f"""
        <div style="background:#F0F4FF;border-left:4px solid #1A2B4A;padding:0.75rem 1rem;border-radius:0 6px 6px 0;font-size:0.9rem;color:#1A2B4A;">
            Ready to generate: <strong>{client_name} — {lob_label}</strong><br>
            <span style="color:#8C8C8C;font-size:0.8rem;">{len(dff):,} rows | {fmt_currency(kpis['spend'])} spend | {fmt_int(kpis['conversions'])} conversions</span>
        </div>
        """, unsafe_allow_html=True)

    with gen_col2:
        if st.button("⚡ Generate Performance Deck", use_container_width=True):
            with st.spinner("Building your deck..."):
                config = {
                    "lob": lob_label,
                    "client": client_name,
                    "date_range": date_range_label,
                }
                buf, err = generate_deck(dff, config)
            if err:
                st.error(f"Error: {err}")
            elif buf:
                st.session_state["deck_buf"] = buf
                st.session_state["deck_name"] = f"{client_name}_{lob_label.replace(' ','_')}_Performance.pptx"
                st.success("Deck generated!")

    if "deck_buf" in st.session_state:
        st.download_button(
            label="⬇ Download PPTX",
            data=st.session_state["deck_buf"],
            file_name=st.session_state.get("deck_name", "performance_deck.pptx"),
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
